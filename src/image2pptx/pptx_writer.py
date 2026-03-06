from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Emu, Inches, Pt

from .config import ExportSettings
from .scene_graph import FreeformNode, PictureAssetNode, PrimitiveNode, SceneGraph, SvgAssetNode

EMU_PER_INCH = 914400


class PptxWriter:
    def __init__(self, export_settings: ExportSettings) -> None:
        self.export_settings = export_settings

    def write(self, scene_graph: SceneGraph, output_path: Path) -> Path:
        presentation = Presentation()
        presentation.slide_width = Inches(self.export_settings.slide_width_in)
        if self.export_settings.slide_height_in is not None:
            presentation.slide_height = Inches(self.export_settings.slide_height_in)
        else:
            presentation.slide_height = Emu(
                int(round(presentation.slide_width * scene_graph.canvas_height / scene_graph.canvas_width))
            )

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*scene_graph.background_color[:3])
        pixels_per_inch = scene_graph.canvas_width / self.export_settings.slide_width_in

        for node in sorted(scene_graph.nodes, key=lambda item: item.z_index):
            if isinstance(node, PrimitiveNode):
                self._add_primitive(slide, node, pixels_per_inch)
            elif isinstance(node, FreeformNode):
                self._add_freeform(slide, node, pixels_per_inch)
            elif isinstance(node, SvgAssetNode):
                self._add_svg_asset(slide, node, pixels_per_inch)
            elif isinstance(node, PictureAssetNode):
                self._add_picture_asset(slide, node, pixels_per_inch)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))
        return output_path

    def _add_primitive(self, slide, node: PrimitiveNode, pixels_per_inch: float) -> None:
        if node.primitive_type == "line" and node.start is not None and node.end is not None:
            shape = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                self._px_to_emu(node.start.x, pixels_per_inch),
                self._px_to_emu(node.start.y, pixels_per_inch),
                self._px_to_emu(node.end.x, pixels_per_inch),
                self._px_to_emu(node.end.y, pixels_per_inch),
            )
            self._style_line(shape, node.stroke_color or node.fill_color, max(1.0, node.stroke_width), pixels_per_inch)
            return

        if node.primitive_type == "circle":
            shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
        elif node.primitive_type == "text":
            shape = slide.shapes.add_textbox(
                self._px_to_emu(node.bbox.x, pixels_per_inch),
                self._px_to_emu(node.bbox.y, pixels_per_inch),
                self._px_to_emu(node.bbox.width, pixels_per_inch),
                self._px_to_emu(node.bbox.height, pixels_per_inch),
            )
            self._style_text(shape, node, pixels_per_inch)
            return
        else:
            shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE

        shape = slide.shapes.add_shape(
            shape_type,
            self._px_to_emu(node.bbox.x, pixels_per_inch),
            self._px_to_emu(node.bbox.y, pixels_per_inch),
            self._px_to_emu(node.bbox.width, pixels_per_inch),
            self._px_to_emu(node.bbox.height, pixels_per_inch),
        )
        self._style_fill(shape, node.fill_color)
        self._style_line(shape, node.stroke_color, node.stroke_width, pixels_per_inch)

    def _add_freeform(self, slide, node: FreeformNode, pixels_per_inch: float) -> None:
        if len(node.points) < 2:
            return

        first = node.points[0]
        builder = slide.shapes.build_freeform(
            start_x=first.x,
            start_y=first.y,
            scale=pixels_per_inch,
        )
        builder.add_line_segments([(point.x, point.y) for point in node.points[1:]], close=node.closed)
        shape = builder.convert_to_shape()
        self._style_fill(shape, node.fill_color)
        self._style_line(shape, node.stroke_color, node.stroke_width, pixels_per_inch)

    def _add_svg_asset(self, slide, node: SvgAssetNode, pixels_per_inch: float) -> None:
        if node.fallback_image_path:
            slide.shapes.add_picture(
                node.fallback_image_path,
                self._px_to_emu(node.bbox.x, pixels_per_inch),
                self._px_to_emu(node.bbox.y, pixels_per_inch),
                self._px_to_emu(node.bbox.width, pixels_per_inch),
                self._px_to_emu(node.bbox.height, pixels_per_inch),
            )
            return

        if node.fallback_points:
            freeform_node = FreeformNode(
                id=node.id,
                bbox=node.bbox,
                z_index=node.z_index,
                fill_color=node.fill_color,
                stroke_color=node.stroke_color,
                stroke_width=node.stroke_width,
                opacity=node.opacity,
                points=node.fallback_points,
                closed=True,
                source_region_id=node.source_region_id,
            )
            self._add_freeform(slide, freeform_node, pixels_per_inch)

    def _add_picture_asset(self, slide, node: PictureAssetNode, pixels_per_inch: float) -> None:
        slide.shapes.add_picture(
            node.image_path,
            self._px_to_emu(node.bbox.x, pixels_per_inch),
            self._px_to_emu(node.bbox.y, pixels_per_inch),
            self._px_to_emu(node.bbox.width, pixels_per_inch),
            self._px_to_emu(node.bbox.height, pixels_per_inch),
        )

    def _style_fill(self, shape, color) -> None:
        if color is None:
            shape.fill.background()
            return
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color[:3])
        if len(color) > 3:
            shape.fill.transparency = max(0.0, min(1.0, 1 - (color[3] / 255)))

    def _style_line(self, shape, color, stroke_width: float, pixels_per_inch: float) -> None:
        if color is None or stroke_width <= 0:
            shape.line.fill.background()
            return
        shape.line.color.rgb = RGBColor(*color[:3])
        shape.line.width = Emu(self._px_to_emu(stroke_width, pixels_per_inch))

    def _style_text(self, shape, node: PrimitiveNode, pixels_per_inch: float) -> None:
        shape.fill.background()
        shape.line.fill.background()
        if not node.text:
            return
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = node.text
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(max(8.0, (node.font_size or node.bbox.height) / pixels_per_inch * 72.0))
        color = node.text_color or (0, 0, 0, 255)
        run.font.color.rgb = RGBColor(*color[:3])

    def _px_to_emu(self, value: float, pixels_per_inch: float) -> int:
        return int(round((value / pixels_per_inch) * EMU_PER_INCH))
